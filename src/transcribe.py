from fastapi import FastAPI, UploadFile, File, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import whisper
import tempfile
import os
import uvicorn
from typing import Optional, List
import logging
import asyncio
import aiofiles
from pathlib import Path
import subprocess
import shutil
from datetime import datetime
import hashlib
import json

# OCR e processamento de documentos
import pytesseract
import easyocr
from PIL import Image
from pdf2image import convert_from_path
import fitz  # PyMuPDF
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
import io

# Configurar logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(
    title="Universal Transcription API",
    description="API universal para transcrição de áudio, OCR de imagens e extração de texto de documentos",
    version="2.0.0"
)

# Configurar CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configurações otimizadas
CACHE_DIR = Path("./cache")
CACHE_DIR.mkdir(exist_ok=True)

MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB
CHUNK_SIZE = 1024 * 1024  # 1MB chunks para upload

# Cache de resultados
transcription_cache = {}

# Carregar modelos Whisper otimizados
logger.info("Carregando modelos Whisper...")
models = {
    "tiny": whisper.load_model("tiny"),    # Para áudios curtos do WhatsApp
    "base": whisper.load_model("base"),    # Para uso geral
    "small": whisper.load_model("small")   # Para vídeos maiores
}
logger.info("Modelos Whisper carregados com sucesso!")

# Inicializar OCRs
logger.info("Inicializando OCRs...")
try:
    # EasyOCR - melhor para textos complexos
    easyocr_reader = easyocr.Reader(['pt', 'en'], gpu=False)
    logger.info("EasyOCR inicializado com sucesso!")
except Exception as e:
    logger.warning(f"Falha ao inicializar EasyOCR: {e}")
    easyocr_reader = None

# Configurar Tesseract (se disponível)
try:
    # Tentar detectar Tesseract
    pytesseract.pytesseract.tesseract_cmd = 'tesseract'  # Linux/Docker
    pytesseract.get_tesseract_version()
    tesseract_available = True
    logger.info("Tesseract OCR disponível!")
except:
    tesseract_available = False
    logger.warning("Tesseract OCR não encontrado")

def get_file_hash(content: bytes) -> str:
    """Gera hash do arquivo para cache"""
    return hashlib.md5(content).hexdigest()

def choose_optimal_model(file_size: int, duration: float = None) -> str:
    """Escolhe o modelo ideal baseado no tamanho e duração"""
    if file_size < 1024 * 1024:  # < 1MB (típico do WhatsApp)
        return "tiny"
    elif file_size < 10 * 1024 * 1024:  # < 10MB
        return "base"
    else:
        return "small"

async def convert_video_to_audio(video_path: str) -> str:
    """Converte vídeo para áudio usando ffmpeg"""
    audio_path = video_path.rsplit('.', 1)[0] + '_audio.wav'
    
    cmd = [
        'ffmpeg', '-i', video_path,
        '-vn',  # Sem vídeo
        '-acodec', 'pcm_s16le',  # Codec de áudio
        '-ar', '16000',  # Sample rate otimizado para Whisper
        '-ac', '1',  # Mono
        '-y',  # Sobrescrever
        audio_path
    ]
    
    try:
        result = subprocess.run(cmd, capture_output=True, text=True, check=True)
        return audio_path
    except subprocess.CalledProcessError as e:
        logger.error(f"Erro na conversão: {e.stderr}")
        raise HTTPException(status_code=500, detail="Erro na conversão do vídeo")

def preprocess_whatsapp_audio(audio_path: str) -> str:
    """Otimiza áudio do WhatsApp para transcrição"""
    output_path = audio_path.rsplit('.', 1)[0] + '_processed.wav'
    
    cmd = [
        'ffmpeg', '-i', audio_path,
        '-ar', '16000',  # Sample rate ideal
        '-ac', '1',      # Mono
        '-filter:a', 'volume=2.0,highpass=f=200,lowpass=f=3000',  # Filtros para voz
        '-y',
        output_path
    ]
    
    try:
        subprocess.run(cmd, capture_output=True, check=True)
        return output_path
    except subprocess.CalledProcessError:
        logger.warning("Falha no pré-processamento, usando arquivo original")
        return audio_path

# ========== FUNÇÕES DE OCR E PROCESSAMENTO DE DOCUMENTOS ==========

def extract_text_from_image(image_path: str, method: str = "easyocr") -> dict:
    """Extrai texto de imagem usando OCR"""
    try:
        if method == "easyocr" and easyocr_reader:
            # EasyOCR - melhor para textos complexos
            results = easyocr_reader.readtext(image_path)
            text = " ".join([result[1] for result in results])
            confidence = sum([result[2] for result in results]) / len(results) if results else 0
            
            return {
                "text": text.strip(),
                "method": "EasyOCR",
                "confidence": confidence,
                "details": results
            }
            
        elif method == "tesseract" and tesseract_available:
            # Tesseract OCR
            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang='por+eng')
            
            # Obter dados detalhados
            data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT, lang='por+eng')
            confidences = [int(conf) for conf in data['conf'] if int(conf) > 0]
            avg_confidence = sum(confidences) / len(confidences) if confidences else 0
            
            return {
                "text": text.strip(),
                "method": "Tesseract",
                "confidence": avg_confidence / 100,
                "details": data
            }
        else:
            raise Exception("Nenhum OCR disponível")
            
    except Exception as e:
        logger.error(f"Erro no OCR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro no OCR: {str(e)}")

def extract_text_from_pdf(pdf_path: str, method: str = "auto") -> dict:
    """Extrai texto de PDF"""
    try:
        text_content = ""
        pages_processed = 0
        
        if method in ["auto", "direct"]:
            # Tentar extrair texto direto primeiro (PDFs com texto)
            try:
                with pdfplumber.open(pdf_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text_content += page_text + "\n"
                            pages_processed += 1
                
                if text_content.strip():
                    return {
                        "text": text_content.strip(),
                        "method": "Direct extraction",
                        "pages": pages_processed,
                        "type": "text_pdf"
                    }
            except Exception as e:
                logger.warning(f"Falha na extração direta: {e}")
        
        if method in ["auto", "ocr"] and not text_content.strip():
            # PDF escaneado - usar OCR
            logger.info("PDF parece ser escaneado, usando OCR...")
            
            # Converter PDF para imagens
            images = convert_from_path(pdf_path, dpi=200)
            
            for i, image in enumerate(images):
                # Salvar imagem temporária
                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_img:
                    image.save(temp_img.name, 'PNG')
                    
                    # Extrair texto da imagem
                    result = extract_text_from_image(temp_img.name)
                    text_content += result["text"] + "\n"
                    pages_processed += 1
                    
                    # Limpar arquivo temporário
                    os.unlink(temp_img.name)
            
            return {
                "text": text_content.strip(),
                "method": "OCR (scanned PDF)",
                "pages": pages_processed,
                "type": "scanned_pdf"
            }
        
        if not text_content.strip():
            raise Exception("Não foi possível extrair texto do PDF")
            
    except Exception as e:
        logger.error(f"Erro ao processar PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro ao processar PDF: {str(e)}")

def extract_text_from_docx(docx_path: str) -> dict:
    """Extrai texto de arquivo Word"""
    try:
        doc = Document(docx_path)
        text_content = ""
        
        for paragraph in doc.paragraphs:
            text_content += paragraph.text + "\n"
        
        # Extrair texto de tabelas também
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text_content += cell.text + " "
                text_content += "\n"
        
        return {
            "text": text_content.strip(),
            "method": "Direct extraction",
            "paragraphs": len(doc.paragraphs),
            "tables": len(doc.tables)
        }
        
    except Exception as e:
        logger.error(f"Erro ao processar DOCX: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro ao processar DOCX: {str(e)}")

def extract_text_from_excel(excel_path: str) -> dict:
    """Extrai texto de arquivo Excel"""
    try:
        workbook = openpyxl.load_workbook(excel_path, data_only=True)
        text_content = ""
        sheets_processed = 0
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_content += f"\n=== {sheet_name} ===\n"
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text_content += row_text + "\n"
            
            sheets_processed += 1
        
        return {
            "text": text_content.strip(),
            "method": "Direct extraction",
            "sheets": sheets_processed
        }
        
    except Exception as e:
        logger.error(f"Erro ao processar Excel: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro ao processar Excel: {str(e)}")

def extract_text_from_pptx(pptx_path: str) -> dict:
    """Extrai texto de arquivo PowerPoint"""
    try:
        presentation = Presentation(pptx_path)
        text_content = ""
        slides_processed = 0
        
        for i, slide in enumerate(presentation.slides):
            text_content += f"\n=== Slide {i+1} ===\n"
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content += shape.text + "\n"
            
            slides_processed += 1
        
        return {
            "text": text_content.strip(),
            "method": "Direct extraction",
            "slides": slides_processed
        }
        
    except Exception as e:
        logger.error(f"Erro ao processar PPTX: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro ao processar PPTX: {str(e)}")

@app.get("/")
async def root():
    return {"message": "Universal Transcription API está funcionando!"}

@app.get("/health")
async def health_check():
    ocr_status = {
        "easyocr": easyocr_reader is not None,
        "tesseract": tesseract_available
    }
    
    return {
        "status": "healthy", 
        "services": {
            "whisper": {
                "available": True,
                "models": list(models.keys())
            },
            "ocr": ocr_status,
            "document_processing": True
        },
        "cache_size": len(transcription_cache),
        "max_file_size_mb": MAX_FILE_SIZE // (1024 * 1024)
    }

@app.post("/transcribe")
async def transcribe_audio(
    file: UploadFile = File(...),
    language: Optional[str] = "pt",
    use_cache: bool = True,
    whatsapp_optimization: bool = False
):
    """
    Transcreve arquivo de áudio/vídeo para texto com otimizações
    
    - **file**: Arquivo de áudio/vídeo
    - **language**: Código do idioma (pt, en, es, etc.)
    - **use_cache**: Usar cache de resultados
    - **whatsapp_optimization**: Aplicar filtros específicos para WhatsApp
    """
    
    # Verificar tamanho do arquivo
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"Arquivo muito grande. Máximo: {MAX_FILE_SIZE // (1024*1024)}MB"
        )
    
    # Verificar cache
    file_hash = get_file_hash(content)
    cache_key = f"{file_hash}_{language}_{whatsapp_optimization}"
    
    if use_cache and cache_key in transcription_cache:
        logger.info(f"Resultado encontrado no cache para {file.filename}")
        return transcription_cache[cache_key]
    
    # Verificar tipo de arquivo
    allowed_types = {
        'audio/mpeg', 'audio/wav', 'audio/mp4', 'audio/m4a', 
        'audio/ogg', 'audio/webm', 'audio/flac',
        'video/mp4', 'video/avi', 'video/mov', 'video/mkv'
    }
    
    if file.content_type not in allowed_types:
        # Tentar detectar pelo nome do arquivo
        ext = file.filename.split('.')[-1].lower() if file.filename else ""
        if ext not in ['mp3', 'wav', 'm4a', 'ogg', 'webm', 'flac', 'mp4', 'avi', 'mov', 'mkv']:
        raise HTTPException(
            status_code=400, 
                detail=f"Tipo de arquivo não suportado: {file.content_type}"
        )
    
    temp_files = []
    try:
        # Criar arquivo temporário
        suffix = f".{file.filename.split('.')[-1]}" if file.filename else ".tmp"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name
            temp_files.append(temp_file_path)
        
        logger.info(f"Processando: {file.filename} ({len(content)} bytes)")
        
        # Processar vídeo se necessário
        is_video = any(ext in file.content_type for ext in ['video/', '.mp4', '.avi', '.mov', '.mkv'])
        if is_video:
            logger.info("Convertendo vídeo para áudio...")
            audio_path = await convert_video_to_audio(temp_file_path)
            temp_files.append(audio_path)
            temp_file_path = audio_path
        
        # Otimizar áudio do WhatsApp
        if whatsapp_optimization:
            logger.info("Aplicando otimizações para WhatsApp...")
            processed_path = preprocess_whatsapp_audio(temp_file_path)
            if processed_path != temp_file_path:
                temp_files.append(processed_path)
                temp_file_path = processed_path
        
        # Escolher modelo otimizado
        model_name = choose_optimal_model(len(content))
        model = models[model_name]
        
        logger.info(f"Usando modelo: {model_name}")
        
        # Transcrever
        transcribe_options = {
            "fp16": False,  # Melhor compatibilidade
            "temperature": 0.0,  # Mais determinístico
        }
        
        if language and language != "auto":
            transcribe_options["language"] = language
        
        result = model.transcribe(temp_file_path, **transcribe_options)
        
        # Preparar resposta
        response = {
            "text": result["text"].strip(),
            "language": result["language"],
            "segments": result["segments"],
            "filename": file.filename,
            "model_used": model_name,
            "duration": result.get("segments", [])[-1]["end"] if result.get("segments") else 0,
            "cached": False
        }
        
        # Salvar no cache
        if use_cache:
            transcription_cache[cache_key] = response.copy()
            transcription_cache[cache_key]["cached"] = True
        
        logger.info("Transcrição concluída com sucesso")
        return response
        
    except Exception as e:
        logger.error(f"Erro na transcrição: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro na transcrição: {str(e)}")
    
    finally:
        # Limpar arquivos temporários
        for temp_path in temp_files:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Erro ao limpar arquivo {temp_path}: {e}")

@app.post("/transcribe-simple")
async def transcribe_simple(
    file: UploadFile = File(...),
    whatsapp_optimization: bool = True
):
    """
    Versão simples otimizada para WhatsApp - retorna apenas o texto
    """
    result = await transcribe_audio(
        file=file, 
        language="pt",
        use_cache=True,
        whatsapp_optimization=whatsapp_optimization
    )
    
    return {"text": result["text"]}

@app.post("/transcribe-whatsapp")
async def transcribe_whatsapp(file: UploadFile = File(...)):
    """
    Endpoint específico para áudios do WhatsApp com otimizações automáticas
    """
    result = await transcribe_audio(
        file=file,
        language="pt", 
        use_cache=True,
        whatsapp_optimization=True
    )
    
    return {
        "text": result["text"],
        "duration": result["duration"],
        "model_used": result["model_used"],
        "cached": result["cached"]
    }

@app.post("/transcribe-video")
async def transcribe_video(
    file: UploadFile = File(...),
    language: Optional[str] = "pt"
):
    """
    Endpoint específico para vídeos grandes
    """
    return await transcribe_audio(
        file=file,
        language=language,
        use_cache=True,
        whatsapp_optimization=False
    )

@app.get("/cache/clear")
async def clear_cache():
    """
    Limpa o cache de transcrições
    """
    global transcription_cache
    cache_size = len(transcription_cache)
    transcription_cache.clear()
    return {"message": f"Cache limpo. {cache_size} itens removidos."}

@app.get("/cache/stats")
async def cache_stats():
    """
    Estatísticas do cache
    """
    return {
        "cache_size": len(transcription_cache),
        "memory_usage_mb": sum(len(str(v)) for v in transcription_cache.values()) / (1024 * 1024)
    }

# ========== ENDPOINTS DE OCR E PROCESSAMENTO DE DOCUMENTOS ==========

@app.post("/ocr/image")
async def ocr_image(
    file: UploadFile = File(...),
    method: str = "easyocr",
    use_cache: bool = True
):
    """
    Extrai texto de imagem usando OCR
    
    - **file**: Arquivo de imagem (jpg, png, bmp, tiff, etc.)
    - **method**: Método OCR ("easyocr" ou "tesseract")
    - **use_cache**: Usar cache de resultados
    """
    
    # Verificar tipo de arquivo
    allowed_types = {
        'image/jpeg', 'image/jpg', 'image/png', 'image/bmp', 
        'image/tiff', 'image/webp', 'image/gif'
    }
    
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"Arquivo muito grande. Máximo: {MAX_FILE_SIZE // (1024*1024)}MB"
        )
    
    # Verificar cache
    file_hash = get_file_hash(content)
    cache_key = f"ocr_{file_hash}_{method}"
    
    if use_cache and cache_key in transcription_cache:
        logger.info(f"OCR encontrado no cache para {file.filename}")
        return transcription_cache[cache_key]
    
    # Verificar extensão se content_type falhar
    if file.content_type not in allowed_types:
        ext = file.filename.split('.')[-1].lower() if file.filename else ""
        if ext not in ['jpg', 'jpeg', 'png', 'bmp', 'tiff', 'webp', 'gif']:
            raise HTTPException(
                status_code=400, 
                detail=f"Tipo de arquivo não suportado: {file.content_type}"
            )
    
    temp_files = []
    try:
        # Criar arquivo temporário
        suffix = f".{file.filename.split('.')[-1]}" if file.filename else ".jpg"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name
            temp_files.append(temp_file_path)
        
        logger.info(f"Processando OCR: {file.filename} ({len(content)} bytes)")
        
        # Extrair texto
        result = extract_text_from_image(temp_file_path, method)
        
        # Preparar resposta
        response = {
            **result,
            "filename": file.filename,
            "file_size": len(content),
            "cached": False
        }
        
        # Salvar no cache
        if use_cache:
            transcription_cache[cache_key] = response.copy()
            transcription_cache[cache_key]["cached"] = True
        
        logger.info("OCR concluído com sucesso")
        return response
        
    except Exception as e:
        logger.error(f"Erro no OCR: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro no OCR: {str(e)}")
    
    finally:
        # Limpar arquivos temporários
        for temp_path in temp_files:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Erro ao limpar arquivo {temp_path}: {e}")

@app.post("/extract/pdf")
async def extract_pdf(
    file: UploadFile = File(...),
    method: str = "auto",
    use_cache: bool = True
):
    """
    Extrai texto de arquivo PDF
    
    - **file**: Arquivo PDF
    - **method**: Método de extração ("auto", "direct", "ocr")
    - **use_cache**: Usar cache de resultados
    """
    
    content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"Arquivo muito grande. Máximo: {MAX_FILE_SIZE // (1024*1024)}MB"
        )
    
    # Verificar cache
    file_hash = get_file_hash(content)
    cache_key = f"pdf_{file_hash}_{method}"
    
    if use_cache and cache_key in transcription_cache:
        logger.info(f"PDF encontrado no cache para {file.filename}")
        return transcription_cache[cache_key]
    
    if not file.filename.lower().endswith('.pdf'):
        raise HTTPException(status_code=400, detail="Arquivo deve ser um PDF")
    
    temp_files = []
    try:
        # Criar arquivo temporário
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name
            temp_files.append(temp_file_path)
        
        logger.info(f"Processando PDF: {file.filename} ({len(content)} bytes)")
        
        # Extrair texto
        result = extract_text_from_pdf(temp_file_path, method)
        
        # Preparar resposta
        response = {
            **result,
            "filename": file.filename,
            "file_size": len(content),
            "cached": False
        }
        
        # Salvar no cache
        if use_cache:
            transcription_cache[cache_key] = response.copy()
            transcription_cache[cache_key]["cached"] = True
        
        logger.info("Extração de PDF concluída com sucesso")
        return response
        
    except Exception as e:
        logger.error(f"Erro ao processar PDF: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro ao processar PDF: {str(e)}")
    
    finally:
        # Limpar arquivos temporários
        for temp_path in temp_files:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Erro ao limpar arquivo {temp_path}: {e}")

@app.post("/extract/document")
async def extract_document(
    file: UploadFile = File(...),
    use_cache: bool = True
):
    """
    Extrai texto de documentos Office (Word, Excel, PowerPoint)
    
    - **file**: Arquivo de documento (.docx, .xlsx, .pptx)
    - **use_cache**: Usar cache de resultados
    """
    
            content = await file.read()
    if len(content) > MAX_FILE_SIZE:
        raise HTTPException(
            status_code=413,
            detail=f"Arquivo muito grande. Máximo: {MAX_FILE_SIZE // (1024*1024)}MB"
        )
    
    # Verificar cache
    file_hash = get_file_hash(content)
    cache_key = f"doc_{file_hash}"
    
    if use_cache and cache_key in transcription_cache:
        logger.info(f"Documento encontrado no cache para {file.filename}")
        return transcription_cache[cache_key]
    
    # Determinar tipo de documento
    filename = file.filename.lower() if file.filename else ""
    
    if not any(filename.endswith(ext) for ext in ['.docx', '.xlsx', '.pptx']):
        raise HTTPException(
            status_code=400, 
            detail="Arquivo deve ser .docx, .xlsx ou .pptx"
        )
    
    temp_files = []
    try:
        # Criar arquivo temporário
        suffix = f".{filename.split('.')[-1]}"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as temp_file:
            temp_file.write(content)
            temp_file_path = temp_file.name
            temp_files.append(temp_file_path)
        
        logger.info(f"Processando documento: {file.filename} ({len(content)} bytes)")
        
        # Extrair texto baseado no tipo
        if filename.endswith('.docx'):
            result = extract_text_from_docx(temp_file_path)
        elif filename.endswith('.xlsx'):
            result = extract_text_from_excel(temp_file_path)
        elif filename.endswith('.pptx'):
            result = extract_text_from_pptx(temp_file_path)
        else:
            raise HTTPException(status_code=400, detail="Tipo de documento não suportado")
        
        # Preparar resposta
        response = {
            **result,
            "filename": file.filename,
            "file_size": len(content),
            "document_type": suffix[1:],
            "cached": False
        }
        
        # Salvar no cache
        if use_cache:
            transcription_cache[cache_key] = response.copy()
            transcription_cache[cache_key]["cached"] = True
        
        logger.info("Extração de documento concluída com sucesso")
        return response
        
    except Exception as e:
        logger.error(f"Erro ao processar documento: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Erro ao processar documento: {str(e)}")
    
    finally:
        # Limpar arquivos temporários
        for temp_path in temp_files:
            try:
                if os.path.exists(temp_path):
                    os.unlink(temp_path)
            except Exception as e:
                logger.warning(f"Erro ao limpar arquivo {temp_path}: {e}")

@app.post("/extract/auto")
async def extract_auto(
    file: UploadFile = File(...),
    use_cache: bool = True
):
    """
    Detecção automática do tipo de arquivo e extração de texto
    
    - **file**: Qualquer arquivo suportado
    - **use_cache**: Usar cache de resultados
    """
    
    if not file.filename:
        raise HTTPException(status_code=400, detail="Nome do arquivo é obrigatório")
    
    filename = file.filename.lower()
    
    # Redirecionar para endpoint apropriado baseado na extensão
    if any(filename.endswith(ext) for ext in ['.mp3', '.wav', '.ogg', '.m4a', '.flac']):
        return await transcribe_whatsapp(file)
    
    elif any(filename.endswith(ext) for ext in ['.mp4', '.avi', '.mov', '.mkv']):
        return await transcribe_video(file)
    
    elif any(filename.endswith(ext) for ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.webp', '.gif']):
        return await ocr_image(file, "easyocr", use_cache)
    
    elif filename.endswith('.pdf'):
        return await extract_pdf(file, "auto", use_cache)
    
    elif any(filename.endswith(ext) for ext in ['.docx', '.xlsx', '.pptx']):
        return await extract_document(file, use_cache)
    
    else:
        raise HTTPException(
            status_code=400, 
            detail=f"Tipo de arquivo não suportado: {filename.split('.')[-1] if '.' in filename else 'desconhecido'}"
        )

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)